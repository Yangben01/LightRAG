"""
此模块包含LightRAG API的所有文档相关路由。
"""

import asyncio
from functools import lru_cache
from lightrag.utils import logger, get_pinyin_sort_key
import aiofiles
import shutil
import traceback
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Any, Literal
from io import BytesIO
from fastapi import (
    APIRouter,
    BackgroundTasks,
    Depends,
    File,
    HTTPException,
    UploadFile,
)
from pydantic import BaseModel, Field, field_validator

from lightrag import LightRAG
from lightrag.base import DeletionResult, DocProcessingStatus, DocStatus
from lightrag.utils import (
    generate_track_id,
    compute_mdhash_id,
    sanitize_text_for_encoding,
)
from lightrag.api.utils_api import get_combined_auth_dependency
from ..config import global_args


@lru_cache(maxsize=1)
def _is_docling_available() -> bool:
    """检查docling是否可用（缓存检查）。
此函数使用lru_cache避免重复导入尝试。
    结果在第一次调用后会被缓存。

    返回:
        bool: 如果docling可用则返回True，否则返回False
    """
    try:
        import docling  # noqa: F401  # type: ignore[import-not-found]

        return True
    except ImportError:
        return False


# 将datetime格式化为带时区信息的ISO格式字符串的函数
def format_datetime(dt: Any) -> Optional[str]:
    """将datetime格式化为带时区信息的ISO格式字符串

    参数:
        dt: Datetime对象、字符串或None

    返回:
        带时区信息的ISO格式字符串，如果输入为None则返回None
    """
    if dt is None:
        return None
    if isinstance(dt, str):
        return dt

    # 检查datetime对象是否有时区信息
    if isinstance(dt, datetime):
        # 如果datetime对象没有时区信息（天真的datetime），则添加UTC时区
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)

    # 返回带时区信息的ISO格式字符串
    return dt.isoformat()
# 返回带时区信息的ISO格式字符串
    return dt.isoformat()


router = APIRouter(
    prefix="/documents",
    tags=["文档"],
)

# 临时文件前缀
temp_prefix = "__tmp__"


def sanitize_filename(filename: str, input_dir: Path) -> str:
    """
    清理上传的文件名以防止路径遍历攻击。

    参数:
        filename: 上传文件的原始文件名
        input_dir: 目标输入目录

    返回:
        str: 清理后的安全文件名

    异常:
        HTTPException: 如果文件名不安全或无效
    """
    # 基本验证
    if not filename or not filename.strip():
        raise HTTPException(status_code=400, detail="文件名不能为空")

    # 删除路径分隔符和遍历序列
    clean_name = filename.replace("/", "").replace("\\", "")
    clean_name = clean_name.replace("..", "")

    # 删除控制字符和空字节
    clean_name = "".join(c for c in clean_name if ord(c) >= 32 and c != "\x7f")

    # 删除前导/尾随空白和点
    clean_name = clean_name.strip().strip(".")
# 检查清理后是否还剩下内容
    if not clean_name:
        raise HTTPException(
            status_code=400, detail="清理后的文件名无效"
        )

    # 验证最终路径是否在输入目录内
    try:
        final_path = (input_dir / clean_name).resolve()
        if not final_path.is_relative_to(input_dir.resolve()):
            raise HTTPException(status_code=400, detail="检测到不安全的文件名")
    except (OSError, ValueError):
        raise HTTPException(status_code=400, detail="无效的文件名")

    return clean_name


class ScanResponse(BaseModel):
    """文档扫描操作的响应模型

    属性:
        status: 扫描操作的状态
        message: 包含额外详情的可选消息
        track_id: 用于监控扫描进度的跟踪ID
    """
    status: Literal["scanning_started"] = Field(
        description="扫描操作的状态"
    )
    message: Optional[str] = Field(
        default=None, description="关于扫描操作的额外详情"
    )
    track_id: str = Field(description="用于监控扫描进度的跟踪ID")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "scanning_started",
                "message": "后台扫描进程已启动",
                "track_id": "scan_20250729_170612_abc123",
            }
        }


class ReprocessResponse(BaseModel):
    """重新处理失败文档操作的响应模型

    属性:
        status: 重新处理操作的状态
        message: 描述操作结果的消息
        track_id: 始终为空字符串。重新处理的文档保留其原始track_id。
    """
    status: Literal["reprocessing_started"] = Field(
        description="重新处理操作的状态"
    )
    message: str = Field(description="描述操作的人类可读消息")
    track_id: str = Field(
        default="",
        description="始终为空字符串。重新处理的文档保留其初始上传时的原始track_id。",
    )

    class Config:
        json_schema_extra = {
            "example": {
                "status": "reprocessing_started",
                "message": "后台已启动对失败文档的重新处理",
                "track_id": "",
            }
        }


class CancelPipelineResponse(BaseModel):
    """管道取消操作的响应模型

    属性:
        status: 取消请求的状态
        message: 描述操作结果的消息
    """

    status: Literal["cancellation_requested", "not_busy"] = Field(
        description="取消请求的状态"
    )
    message: str = Field(description="描述操作的人类可读消息")
    
    class Config:
        json_schema_extra = {
            "example": {
                "status": "cancellation_requested",
                "message": "已请求取消管道。文档将被标记为FAILED。",
            }
        }


class InsertTextRequest(BaseModel):
    """插入单个文本文档的请求模型

    属性:
        text: 要插入到RAG系统中的文本内容
        file_source: 文本来源（可选）
    """

    text: str = Field(
        min_length=1,
        description="要插入的文本",
    )
    file_source: str = Field(default=None, min_length=0, description="文件来源")

    @field_validator("text", mode="after")
    @classmethod
    def strip_text_after(cls, text: str) -> str:
        return text.strip()

    @field_validator("file_source", mode="after")
    @classmethod
    def strip_source_after(cls, file_source: str) -> str:
        return file_source.strip()
        
    class Config:
        json_schema_extra = {
            "example": {
                "text": "这是要插入到RAG系统中的示例文本。",
                "file_source": "文本来源（可选）",
            }
        }


class InsertTextsRequest(BaseModel):
    """插入多个文本文档的请求模型

    属性:
        texts: 要插入到RAG系统中的文本内容列表
        file_sources: 文本来源（可选）
    """

    texts: list[str] = Field(
        min_length=1,
        description="要插入的文本",
    )
    file_sources: list[str] = Field(
        default=None, min_length=0, description="文本来源"
    )

    @field_validator("texts", mode="after")
    @classmethod
    def strip_texts_after(cls, texts: list[str]) -> list[str]:
        return [text.strip() for text in texts]
        
    @field_validator("file_sources", mode="after")
    @classmethod
    def strip_sources_after(cls, file_sources: list[str]) -> list[str]:
        return [file_source.strip() for file_source in file_sources]

    class Config:
        json_schema_extra = {
            "example": {
                "texts": [
                    "这是要插入的第一个文本。",
                    "这是要插入的第二个文本。",
                ],
                "file_sources": [
                    "第一个文件来源（可选）",
                ],
            }
        }


class InsertResponse(BaseModel):
    """文档插入操作的响应模型

    属性:
        status: 操作状态（success, duplicated, partial_success, failure）
        message: 描述操作结果的详细消息
        track_id: 用于监控处理状态的跟踪ID
    """
    status: Literal["success", "duplicated", "partial_success", "failure"] = Field(
        description="操作状态"
    )
    message: str = Field(description="描述操作结果的消息")
    track_id: str = Field(description="用于监控处理状态的跟踪ID")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "文件'document.pdf'上传成功。将在后台继续处理。",
                "track_id": "upload_20250729_170612_abc123",
            }
        }


class ClearDocumentsResponse(BaseModel):
    """清除文档操作的响应模型

    属性:
        status: 清除操作的状态
        message: 描述操作结果的详细消息
    """

    status: Literal["success", "partial_success", "busy", "fail"] = Field(
        description="清除操作的状态"
    )
    message: str = Field(description="描述操作结果的消息")
    
    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "所有文档已成功清除。删除了15个文件。",
            }
        }


class ClearCacheRequest(BaseModel):
    """清除缓存的请求模型

    此模型保留用于API兼容性，但不再接受任何参数。
    无论请求内容如何，所有缓存都将被清除。
    """

    class Config:
        json_schema_extra = {"example": {}}


class ClearCacheResponse(BaseModel):
    """清除缓存操作的响应模型

    属性:
        status: 清除操作的状态
        message: 描述操作结果的详细消息
    """

    status: Literal["success", "fail"] = Field(
        description="清除操作的状态"
    )
    message: str = Field(description="描述操作结果的消息")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "已成功清除以下模式的缓存: ['default', 'naive']",
            }
        }


class DeleteDocRequest(BaseModel):
    doc_ids: List[str] = Field(..., description="要删除的文档ID列表。")
    delete_file: bool = Field(
        default=False,
        description="是否删除上传目录中的对应文件。",
    )
    delete_llm_cache: bool = Field(
        default=False,
        description="是否删除文档的缓存LLM提取结果。",
    )
    
    @field_validator("doc_ids", mode="after")
    @classmethod
    def validate_doc_ids(cls, doc_ids: List[str]) -> List[str]:
        if not doc_ids:
            raise ValueError("文档ID列表不能为空")

        validated_ids = []
        for doc_id in doc_ids:
            if not doc_id or not doc_id.strip():
                raise ValueError("文档ID不能为空")
            validated_ids.append(doc_id.strip())

        # 检查重复项
        if len(validated_ids) != len(set(validated_ids)):
            raise ValueError("文档ID必须唯一")

        return validated_ids


class DeleteEntityRequest(BaseModel):
    entity_name: str = Field(..., description="要删除的实体名称。")

    @field_validator("entity_name", mode="after")
    @classmethod
    def validate_entity_name(cls, entity_name: str) -> str:
        if not entity_name or not entity_name.strip():
            raise ValueError("实体名称不能为空")
        return entity_name.strip()


class DeleteRelationRequest(BaseModel):
    source_entity: str = Field(..., description="源实体的名称。")
    target_entity: str = Field(..., description="目标实体的名称。")

    @field_validator("source_entity", "target_entity", mode="after")
    @classmethod
    def validate_entity_names(cls, entity_name: str) -> str:
        if not entity_name or not entity_name.strip():
            raise ValueError("实体名称不能为空")
        return entity_name.strip()


class DocStatusResponse(BaseModel):
    """文档状态的响应模型

    属性:
        id: 文档标识符
        content_summary: 文档内容摘要
        content_length: 文档内容长度
        status: 当前处理状态
        created_at: 创建时间戳（ISO格式字符串）
        updated_at: 最后更新时间戳（ISO格式字符串）
        chunks_count: 块数量（可选）
        error: 错误消息（可选）
        metadata: 附加元数据（可选）
        file_path: 文档文件路径
    """
    id: str = Field(description="文档标识符")
    content_summary: str = Field(description="文档内容摘要")
    content_length: int = Field(description="文档内容的字符长度")
    status: DocStatus = Field(description="当前处理状态")
    created_at: str = Field(description="创建时间戳（ISO格式字符串）")
    updated_at: str = Field(description="最后更新时间戳（ISO格式字符串）")
    track_id: Optional[str] = Field(
        default=None, description="用于监控进度的跟踪ID"
    )
    chunks_count: Optional[int] = Field(
        default=None, description="文档被分割成的块数"
    )
    error_msg: Optional[str] = Field(
        default=None, description="如果处理失败则为错误消息"
    )
    metadata: Optional[dict[str, Any]] = Field(
        default=None, description="关于文档的附加元数据"
    )
    file_path: str = Field(description="文档文件路径")
    
    class Config:
        json_schema_extra = {
            "example": {
                "id": "doc_123456",
                "content_summary": "机器学习研究论文",
                "content_length": 15240,
                "status": "processed",
                "created_at": "2025-03-31T12:34:56",
                "updated_at": "2025-03-31T12:35:30",
                "track_id": "upload_20250729_170612_abc123",
                "chunks_count": 12,
                "error": None,
                "metadata": {"author": "John Doe", "year": 2025},
                "file_path": "research_paper.pdf",
            }
        }


class DocsStatusesResponse(BaseModel):
    """文档状态的响应模型

    属性:
        statuses: 将文档状态映射到文档状态响应列表的字典
    """
    statuses: Dict[DocStatus, List[DocStatusResponse]] = Field(
        default_factory=dict,
        description="将文档状态映射到文档状态响应列表的字典",
    )
    
    class Config:
        json_schema_extra = {
            "example": {
                "statuses": {
                    "PENDING": [
                        {
                            "id": "doc_123",
                            "content_summary": "待处理文档",
                            "content_length": 5000,
                            "status": "pending",
                            "created_at": "2025-03-31T10:00:00",
                            "updated_at": "2025-03-31T10:00:00",
                            "track_id": "upload_20250331_100000_abc123",
                            "chunks_count": None,
                            "error": None,
                            "metadata": None,
                            "file_path": "pending_doc.pdf",
                        }
                    ],
                    "PREPROCESSED": [
                        {
                            "id": "doc_789",
                            "content_summary": "等待最终索引的文档",
                            "content_length": 7200,
                            "status": "preprocessed",
                            "created_at": "2025-03-31T09:30:00",
                            "updated_at": "2025-03-31T09:35:00",
                            "track_id": "upload_20250331_093000_xyz789",
                            "chunks_count": 10,
                            "error": None,
                            "metadata": None,
                            "file_path": "preprocessed_doc.pdf",
                        }
                    ],
                    "PROCESSED": [
                        {
                            "id": "doc_456",
                            "content_summary": "已处理文档",
                            "content_length": 8000,
                            "status": "processed",
                            "created_at": "2025-03-31T09:00:00",
                            "updated_at": "2025-03-31T09:05:00",
                            "track_id": "insert_20250331_090000_def456",
                            "chunks_count": 8,
                            "error": None,
                            "metadata": {"author": "John Doe"},
                            "file_path": "processed_doc.pdf",
                        }
                    ],
                }
            }
        }


class TrackStatusResponse(BaseModel):
    """通过track_id跟踪文档处理状态的响应模型

    属性:
        track_id: 跟踪ID
        documents: 与此track_id关联的文档列表
        total_count: 此track_id的文档总数
        status_summary: 按状态分类的文档计数
    """

    track_id: str = Field(description="跟踪ID")
    documents: List[DocStatusResponse] = Field(
        description="与此track_id关联的文档列表"
    )
    total_count: int = Field(description="此track_id的文档总数")
    status_summary: Dict[str, int] = Field(description="按状态分类的文档计数")
    
    class Config:
        json_schema_extra = {
            "example": {
                "track_id": "upload_20250729_170612_abc123",
                "documents": [
                    {
                        "id": "doc_123456",
                        "content_summary": "机器学习研究论文",
                        "content_length": 15240,
                        "status": "PROCESSED",
                        "created_at": "2025-03-31T12:34:56",
                        "updated_at": "2025-03-31T12:35:30",
                        "track_id": "upload_20250729_170612_abc123",
                        "chunks_count": 12,
                        "error": None,
                        "metadata": {"author": "John Doe", "year": 2025},
                        "file_path": "research_paper.pdf",
                    }
                ],
                "total_count": 1,
                "status_summary": {"PROCESSED": 1},
            }
        }


class DocumentsRequest(BaseModel):
    """分页文档查询的请求模型

    属性:
        status_filter: 按文档状态过滤，None表示所有状态
        page: 页码（从1开始）
        page_size: 每页文档数（10-200）
        sort_field: 排序字段（'created_at', 'updated_at', 'id', 'file_path'）
        sort_direction: 排序方向（'asc' 或 'desc'）
    """
    status_filter: Optional[DocStatus] = Field(
        default=None, description="按文档状态过滤，None表示所有状态"
    )
    page: int = Field(default=1, ge=1, description="页码（从1开始）")
    page_size: int = Field(
        default=50, ge=10, le=200, description="每页文档数（10-200）"
    )
    sort_field: Literal["created_at", "updated_at", "id", "file_path"] = Field(
        default="updated_at", description="排序字段"
    )
    sort_direction: Literal["asc", "desc"] = Field(
        default="desc", description="排序方向"
    )

    class Config:
        json_schema_extra = {
            "example": {
                "status_filter": "PROCESSED",
                "page": 1,
                "page_size": 50,
                "sort_field": "updated_at",
                "sort_direction": "desc",
            }
        }


class PaginationInfo(BaseModel):
    """分页信息

    属性:
        page: 当前页码
        page_size: 每页项目数
        total_count: 项目总数
        total_pages: 总页数
        has_next: 是否有下一页
        has_prev: 是否有上一页
    """

    page: int = Field(description="当前页码")
    page_size: int = Field(description="每页项目数")
    total_count: int = Field(description="项目总数")
    total_pages: int = Field(description="总页数")
    has_next: bool = Field(description="是否有下一页")
    has_prev: bool = Field(description="是否有上一页")

    class Config:
        json_schema_extra = {
            "example": {
                "page": 1,
                "page_size": 50,
                "total_count": 150,
                "total_pages": 3,
                "has_next": True,
                "has_prev": False,
            }
        }


class PaginatedDocsResponse(BaseModel):
    """分页文档查询的响应模型

    属性:
        documents: 当前页面的文档列表
        pagination: 分页信息
        status_counts: 所有文档按状态分类的计数
    """

    documents: List[DocStatusResponse] = Field(
        description="当前页面的文档列表"
    )
    pagination: PaginationInfo = Field(description="分页信息")
    status_counts: Dict[str, int] = Field(
        description="所有文档按状态分类的计数"
    )
    
    class Config:
        json_schema_extra = {
            "example": {
                "documents": [
                    {
                        "id": "doc_123456",
                        "content_summary": "机器学习研究论文",
                        "content_length": 15240,
                        "status": "PROCESSED",
                        "created_at": "2025-03-31T12:34:56",
                        "updated_at": "2025-03-31T12:35:30",
                        "track_id": "upload_20250729_170612_abc123",
                        "chunks_count": 12,
                        "error_msg": None,
                        "metadata": {"author": "John Doe", "year": 2025},
                        "file_path": "research_paper.pdf",
                    }
                ],
                "pagination": {
                    "page": 1,
                    "page_size": 50,
                    "total_count": 150,
                    "total_pages": 3,
                    "has_next": True,
                    "has_prev": False,
                },
                "status_counts": {
                    "PENDING": 10,
                    "PROCESSING": 5,
                    "PREPROCESSED": 5,
                    "PROCESSED": 130,
                    "FAILED": 5,
                },
            }
        }


class StatusCountsResponse(BaseModel):
    """文档状态计数的响应模型

    属性:
        status_counts: 按状态分类的文档计数
    """

    status_counts: Dict[str, int] = Field(description="按状态分类的文档计数")

    class Config:
        json_schema_extra = {
            "example": {
                "status_counts": {
                    "PENDING": 10,
                    "PROCESSING": 5,
                    "PREPROCESSED": 5,
                    "PROCESSED": 130,
                    "FAILED": 5,
                }
            }
        }


class PipelineStatusResponse(BaseModel):
    """管道状态的响应模型
    
    属性:
        autoscanned: 自动扫描是否已启动
        busy: 管道当前是否忙碌
        job_name: 当前作业名称（例如，索引文件/索引文本）
        job_start: 作业开始时间作为带时区的ISO格式字符串（可选）
        docs: 待索引的文档总数
        batchs: 处理文档的批次数
        cur_batch: 当前处理批次
        request_pending: 待处理请求的标志
        latest_message: 管道处理的最新消息
        history_messages: 历史消息列表
        update_status: 所有命名空间的更新标志状态
    """

        docs: Total number of documents to be indexed
        batchs: Number of batches for processing documents
        cur_batch: Current processing batch
        request_pending: Flag for pending request for processing
        latest_message: Latest message from pipeline processing
        history_messages: List of history messages
        update_status: Status of update flags for all namespaces
    """

    autoscanned: bool = False
    busy: bool = False
    job_name: str = "Default Job"
    job_start: Optional[str] = None
    docs: int = 0
    batchs: int = 0
    cur_batch: int = 0
    request_pending: bool = False
    latest_message: str = ""
    history_messages: Optional[List[str]] = None
    update_status: Optional[dict] = None

    @field_validator("job_start", mode="before")
    @classmethod
    def parse_job_start(cls, value):
        """Process datetime and return as ISO format string with timezone"""
        return format_datetime(value)

    class Config:
        extra = "allow"  # Allow additional fields from the pipeline status


class DocumentManager:
    def __init__(
        self,
        input_dir: str,
        workspace: str = "",  # New parameter for workspace isolation
        supported_extensions: tuple = (
            ".txt",
            ".md",
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
            ".rtf",  # Rich Text Format
            ".odt",  # OpenDocument Text
            ".tex",  # LaTeX
            ".epub",  # Electronic Publication
            ".html",  # HyperText Markup Language
            ".htm",  # HyperText Markup Language
            ".csv",  # Comma-Separated Values
            ".json",  # JavaScript Object Notation
            ".xml",  # eXtensible Markup Language
            ".yaml",  # YAML Ain't Markup Language
            ".yml",  # YAML
            ".log",  # Log files
            ".conf",  # Configuration files
            ".ini",  # Initialization files
            ".properties",  # Java properties files
            ".sql",  # SQL scripts
            ".bat",  # Batch files
            ".sh",  # Shell scripts
            ".c",  # C source code
            ".cpp",  # C++ source code
            ".py",  # Python source code
            ".java",  # Java source code
            ".js",  # JavaScript source code
            ".ts",  # TypeScript source code
            ".swift",  # Swift source code
            ".go",  # Go source code
            ".rb",  # Ruby source code
            ".php",  # PHP source code
            ".css",  # Cascading Style Sheets
            ".scss",  # Sassy CSS
            ".less",  # LESS CSS
        ),
    ):
        # Store the base input directory and workspace
        self.base_input_dir = Path(input_dir)
        self.workspace = workspace
        self.supported_extensions = supported_extensions
        self.indexed_files = set()

        # Create workspace-specific input directory
        # If workspace is provided, create a subdirectory for data isolation
        if workspace:
            self.input_dir = self.base_input_dir / workspace
        else:
            self.input_dir = self.base_input_dir

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            logger.debug(f"Scanning for {ext} files in {self.input_dir}")
            for file_path in self.input_dir.glob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)


def validate_file_path_security(file_path_str: str, base_dir: Path) -> Optional[Path]:
    """
    Validate file path security to prevent Path Traversal attacks.

    Args:
        file_path_str: The file path string to validate
        base_dir: The base directory that the file must be within

    Returns:
        Path: Safe file path if valid, None if unsafe or invalid
    """
    if not file_path_str or not file_path_str.strip():
        return None

    try:
        # Clean the file path string
        clean_path_str = file_path_str.strip()

        # Check for obvious path traversal patterns before processing
        # This catches both Unix (..) and Windows (..\) style traversals
        if ".." in clean_path_str:
            # Additional check for Windows-style backslash traversal
            if (
                "\\..\\" in clean_path_str
                or clean_path_str.startswith("..\\")
                or clean_path_str.endswith("\\..")
            ):
                # logger.warning(
                #     f"Security violation: Windows path traversal attempt detected - {file_path_str}"
                # )
                return None

        # Normalize path separators (convert backslashes to forward slashes)
        # This helps handle Windows-style paths on Unix systems
        normalized_path = clean_path_str.replace("\\", "/")

        # Create path object and resolve it (handles symlinks and relative paths)
        candidate_path = (base_dir / normalized_path).resolve()
        base_dir_resolved = base_dir.resolve()

        # Check if the resolved path is within the base directory
        if not candidate_path.is_relative_to(base_dir_resolved):
            # logger.warning(
            #     f"Security violation: Path traversal attempt detected - {file_path_str}"
            # )
            return None

        return candidate_path

    except (OSError, ValueError, Exception) as e:
        logger.warning(f"Invalid file path detected: {file_path_str} - {str(e)}")
        return None


def get_unique_filename_in_enqueued(target_dir: Path, original_name: str) -> str:
    """Generate a unique filename in the target directory by adding numeric suffixes if needed

    Args:
        target_dir: Target directory path
        original_name: Original filename

    Returns:
        str: Unique filename (may have numeric suffix added)
    """
    import time

    original_path = Path(original_name)
    base_name = original_path.stem
    extension = original_path.suffix

    # Try original name first
    if not (target_dir / original_name).exists():
        return original_name

    # Try with numeric suffixes 001-999
    for i in range(1, 1000):
        suffix = f"{i:03d}"
        new_name = f"{base_name}_{suffix}{extension}"
        if not (target_dir / new_name).exists():
            return new_name

    # Fallback with timestamp if all 999 slots are taken
    timestamp = int(time.time())
    return f"{base_name}_{timestamp}{extension}"


# Document processing helper functions (synchronous)
# These functions run in thread pool via asyncio.to_thread() to avoid blocking the event loop


def _convert_with_docling(file_path: Path) -> str:
    """Convert document using docling (synchronous).

    Args:
        file_path: Path to the document file

    Returns:
        str: Extracted markdown content
    """
    from docling.document_converter import DocumentConverter  # type: ignore

    converter = DocumentConverter()
    result = converter.convert(file_path)
    return result.document.export_to_markdown()


def _extract_pdf_pypdf(file_bytes: bytes, password: str = None) -> str:
    """Extract PDF content using pypdf (synchronous).

    Args:
        file_bytes: PDF file content as bytes
        password: Optional password for encrypted PDFs

    Returns:
        str: Extracted text content

    Raises:
        Exception: If PDF is encrypted and password is incorrect or missing
    """
    from pypdf import PdfReader  # type: ignore

    pdf_file = BytesIO(file_bytes)
    reader = PdfReader(pdf_file)

    # Check if PDF is encrypted
    if reader.is_encrypted:
        if not password:
            raise Exception("PDF is encrypted but no password provided")

        decrypt_result = reader.decrypt(password)
        if decrypt_result == 0:
            raise Exception("Incorrect PDF password")

    # Extract text from all pages
    content = ""
    for page in reader.pages:
        content += page.extract_text() + "\n"

    return content


def _extract_docx(file_bytes: bytes) -> str:
    """Extract DOCX content including tables in document order (synchronous).

    Args:
        file_bytes: DOCX file content as bytes

    Returns:
        str: Extracted text content with tables in their original positions.
             Tables are separated from paragraphs with blank lines for clarity.
    """
    from docx import Document  # type: ignore
    from docx.table import Table  # type: ignore
    from docx.text.paragraph import Paragraph  # type: ignore

    docx_file = BytesIO(file_bytes)
    doc = Document(docx_file)

    def escape_cell(cell_value: str | None) -> str:
        """Escape characters that would break tab-delimited layout.

        Escape order is critical: backslashes first, then tabs/newlines.
        This prevents double-escaping issues.

        Args:
            cell_value: The cell value to escape (can be None or str)

        Returns:
            str: Escaped cell value safe for tab-delimited format
        """
        if cell_value is None:
            return ""
        text = str(cell_value)
        # CRITICAL: Escape backslash first to avoid double-escaping
        return (
            text.replace("\\", "\\\\")  # Must be first: \ -> \\
            .replace("\t", "\\t")  # Tab -> \t (visible)
            .replace("\r\n", "\\n")  # Windows newline -> \n
            .replace("\r", "\\n")  # Mac newline -> \n
            .replace("\n", "\\n")  # Unix newline -> \n
        )

    content_parts = []
    in_table = False  # Track if we're currently processing a table

    # Helper function to safely get tag local name
    def get_tag_localname(tag):
        """Safely extract local name from tag, handling QName objects and strings."""
        try:
            # Try to get localname if it's a QName object
            if hasattr(tag, 'localname'):
                return tag.localname
            # Try to convert to string and extract local name
            tag_str = str(tag)
            # Extract local name if tag contains namespace (e.g., "{http://...}p" -> "p")
            if "}" in tag_str:
                return tag_str.split("}")[-1]
            return tag_str
        except Exception:
            # Fallback: try string conversion
            return str(tag).split("}")[-1] if "}" in str(tag) else str(tag)

    # Iterate through all body elements in document order
    for element in doc.element.body:
        tag_local = get_tag_localname(element.tag)
        
        # Check if element is a paragraph
        if tag_local.endswith("p"):
            # If coming out of a table, add blank line after table
            if in_table:
                content_parts.append("")  # Blank line after table
                in_table = False

            paragraph = Paragraph(element, doc)
            text = paragraph.text
            # Always append to preserve document spacing (including blank paragraphs)
            content_parts.append(text)

        # Check if element is a table
        elif tag_local.endswith("tbl"):
            # Add blank line before table (if content exists)
            if content_parts and not in_table:
                content_parts.append("")  # Blank line before table

            in_table = True
            table = Table(element, doc)
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text
                    # Escape special characters to preserve tab-delimited structure
                    row_text.append(escape_cell(cell_text))
                # Only add row if at least one cell has content
                if any(cell for cell in row_text):
                    content_parts.append("\t".join(row_text))

    return "\n".join(content_parts)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract PPTX content (synchronous).

    Args:
        file_bytes: PPTX file content as bytes

    Returns:
        str: Extracted text content
    """
    from pptx import Presentation  # type: ignore

    pptx_file = BytesIO(file_bytes)
    prs = Presentation(pptx_file)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract XLSX content in tab-delimited format with clear sheet separation.

    This function processes Excel workbooks and converts them to a structured text format
    suitable for LLM prompts and RAG systems. Each sheet is clearly delimited with
    separator lines, and special characters are escaped to preserve the tab-delimited structure.

    Features:
    - Each sheet is wrapped with '====================' separators for visual distinction
    - Special characters (tabs, newlines, backslashes) are escaped to prevent structure corruption
    - Column alignment is preserved across all rows to maintain tabular structure
    - Empty rows are preserved as blank lines to maintain row structure
    - Uses sheet.max_column to determine column width efficiently

    Args:
        file_bytes: XLSX file content as bytes

    Returns:
        str: Extracted text content with all sheets in tab-delimited format.
             Format: Sheet separators, sheet name, then tab-delimited rows.

    Example output:
        ==================== Sheet: Data ====================
        Name\tAge\tCity
        Alice\t30\tNew York
        Bob\t25\tLondon

        ==================== Sheet: Summary ====================
        Total\t2
        ====================
    """
    from openpyxl import load_workbook  # type: ignore

    xlsx_file = BytesIO(file_bytes)
    wb = load_workbook(xlsx_file)

    def escape_cell(cell_value: str | int | float | None) -> str:
        """Escape characters that would break tab-delimited layout.

        Escape order is critical: backslashes first, then tabs/newlines.
        This prevents double-escaping issues.

        Args:
            cell_value: The cell value to escape (can be None, str, int, or float)

        Returns:
            str: Escaped cell value safe for tab-delimited format
        """
        if cell_value is None:
            return ""
        text = str(cell_value)
        # CRITICAL: Escape backslash first to avoid double-escaping
        return (
            text.replace("\\", "\\\\")  # Must be first: \ -> \\
            .replace("\t", "\\t")  # Tab -> \t (visible)
            .replace("\r\n", "\\n")  # Windows newline -> \n
            .replace("\r", "\\n")  # Mac newline -> \n
            .replace("\n", "\\n")  # Unix newline -> \n
        )

    def escape_sheet_title(title: str) -> str:
        """Escape sheet title to prevent formatting issues in separators.

        Args:
            title: Original sheet title

        Returns:
            str: Sanitized sheet title with tabs/newlines replaced
        """
        return str(title).replace("\n", " ").replace("\t", " ").replace("\r", " ")

    content_parts: list[str] = []
    sheet_separator = "=" * 20

    for idx, sheet in enumerate(wb):
        if idx > 0:
            content_parts.append("")  # Blank line between sheets for readability

        # Escape sheet title to handle edge cases with special characters
        safe_title = escape_sheet_title(sheet.title)
        content_parts.append(f"{sheet_separator} Sheet: {safe_title} {sheet_separator}")

        # Use sheet.max_column to get the maximum column width directly
        max_columns = sheet.max_column if sheet.max_column else 0

        # Extract rows with consistent width to preserve column alignment
        for row in sheet.iter_rows(values_only=True):
            row_parts = []

            # Build row up to max_columns width
            for idx in range(max_columns):
                if idx < len(row):
                    row_parts.append(escape_cell(row[idx]))
                else:
                    row_parts.append("")  # Pad short rows

            # Check if row is completely empty
            if all(part == "" for part in row_parts):
                # Preserve empty rows as blank lines (maintains row structure)
                content_parts.append("")
            else:
                # Join all columns to maintain consistent column count
                content_parts.append("\t".join(row_parts))

    # Final separator for symmetry (makes parsing easier)
    content_parts.append(sheet_separator)
    return "\n".join(content_parts)


async def pipeline_enqueue_file(
    rag: LightRAG, file_path: Path, track_id: str = None
) -> tuple[bool, str]:
    """Add a file to the queue for processing

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
        track_id: Optional tracking ID, if not provided will be generated
    Returns:
        tuple: (success: bool, track_id: str)
    """

    # Generate track_id if not provided
    if track_id is None:
        track_id = generate_track_id("unknown")

    try:
        content = ""
        ext = file_path.suffix.lower()
        file_size = 0

        # Get file size for error reporting
        try:
            file_size = file_path.stat().st_size
        except Exception:
            file_size = 0

        file = None
        try:
            async with aiofiles.open(file_path, "rb") as f:
                file = await f.read()
        except PermissionError as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]Permission denied - cannot read file",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Permission denied reading file: {file_path.name}"
            )
            return False, track_id
        except FileNotFoundError as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File not found",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(f"[File Extraction]File not found: {file_path.name}")
            return False, track_id
        except Exception as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File reading error",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Error reading file {file_path.name}: {str(e)}"
            )
            return False, track_id

        # Process based on file type
        try:
            match ext:
                case (
                    ".txt"
                    | ".md"
                    | ".html"
                    | ".htm"
                    | ".tex"
                    | ".json"
                    | ".xml"
                    | ".yaml"
                    | ".yml"
                    | ".rtf"
                    | ".odt"
                    | ".epub"
                    | ".csv"
                    | ".log"
                    | ".conf"
                    | ".ini"
                    | ".properties"
                    | ".sql"
                    | ".bat"
                    | ".sh"
                    | ".c"
                    | ".cpp"
                    | ".py"
                    | ".java"
                    | ".js"
                    | ".ts"
                    | ".swift"
                    | ".go"
                    | ".rb"
                    | ".php"
                    | ".css"
                    | ".scss"
                    | ".less"
                ):
                    try:
                        # Try to decode as UTF-8
                        content = file.decode("utf-8")

                        # Validate content
                        if not content or len(content.strip()) == 0:
                            error_files = [
                                {
                                    "file_path": str(file_path.name),
                                    "error_description": "[File Extraction]Empty file content",
                                    "original_error": "File contains no content or only whitespace",
                                    "file_size": file_size,
                                }
                            ]
                            await rag.apipeline_enqueue_error_documents(
                                error_files, track_id
                            )
                            logger.error(
                                f"[File Extraction]Empty content in file: {file_path.name}"
                            )
                            return False, track_id

                        # Check if content looks like binary data string representation
                        if content.startswith("b'") or content.startswith('b"'):
                            error_files = [
                                {
                                    "file_path": str(file_path.name),
                                    "error_description": "[File Extraction]Binary data in text file",
                                    "original_error": "File appears to contain binary data representation instead of text",
                                    "file_size": file_size,
                                }
                            ]
                            await rag.apipeline_enqueue_error_documents(
                                error_files, track_id
                            )
                            logger.error(
                                f"[File Extraction]File {file_path.name} appears to contain binary data representation instead of text"
                            )
                            return False, track_id

                    except UnicodeDecodeError as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]UTF-8 encoding error, please convert it to UTF-8 before processing",
                                "original_error": f"File is not valid UTF-8 encoded text: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]File {file_path.name} is not valid UTF-8 encoded text. Please convert it to UTF-8 before processing."
                        )
                        return False, track_id

                case ".pdf":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to pypdf."
                                )
                            # Use pypdf (non-blocking via to_thread)
                            content = await asyncio.to_thread(
                                _extract_pdf_pypdf,
                                file,
                                global_args.pdf_decrypt_password,
                            )
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]PDF processing error",
                                "original_error": f"Failed to extract text from PDF: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing PDF {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".docx":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to python-docx."
                                )
                            # Use python-docx (non-blocking via to_thread)
                            content = await asyncio.to_thread(_extract_docx, file)
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]DOCX processing error",
                                "original_error": f"Failed to extract text from DOCX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing DOCX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".pptx":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to python-pptx."
                                )
                            # Use python-pptx (non-blocking via to_thread)
                            content = await asyncio.to_thread(_extract_pptx, file)
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]PPTX processing error",
                                "original_error": f"Failed to extract text from PPTX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing PPTX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".xlsx":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to openpyxl."
                                )
                            # Use openpyxl (non-blocking via to_thread)
                            content = await asyncio.to_thread(_extract_xlsx, file)
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]XLSX processing error",
                                "original_error": f"Failed to extract text from XLSX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing XLSX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case _:
                    error_files = [
                        {
                            "file_path": str(file_path.name),
                            "error_description": f"[File Extraction]Unsupported file type: {ext}",
                            "original_error": f"File extension {ext} is not supported",
                            "file_size": file_size,
                        }
                    ]
                    await rag.apipeline_enqueue_error_documents(error_files, track_id)
                    logger.error(
                        f"[File Extraction]Unsupported file type: {file_path.name} (extension {ext})"
                    )
                    return False, track_id

        except Exception as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File format processing error",
                    "original_error": f"Unexpected error during file extracting: {str(e)}",
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Unexpected error during {file_path.name} extracting: {str(e)}"
            )
            return False, track_id

        # Insert into the RAG queue
        if content:
            # Check if content contains only whitespace characters
            if not content.strip():
                error_files = [
                    {
                        "file_path": str(file_path.name),
                        "error_description": "[File Extraction]File contains only whitespace",
                        "original_error": "File content contains only whitespace characters",
                        "file_size": file_size,
                    }
                ]
                await rag.apipeline_enqueue_error_documents(error_files, track_id)
                logger.warning(
                    f"[File Extraction]File contains only whitespace characters: {file_path.name}"
                )
                return False, track_id

            try:
                await rag.apipeline_enqueue_documents(
                    content, file_paths=file_path.name, track_id=track_id
                )

                logger.info(
                    f"Successfully extracted and enqueued file: {file_path.name}"
                )

                # Move file to __enqueued__ directory after enqueuing
                try:
                    enqueued_dir = file_path.parent / "__enqueued__"
                    enqueued_dir.mkdir(exist_ok=True)

                    # Generate unique filename to avoid conflicts
                    unique_filename = get_unique_filename_in_enqueued(
                        enqueued_dir, file_path.name
                    )
                    target_path = enqueued_dir / unique_filename

                    # Move the file
                    file_path.rename(target_path)
                    logger.debug(
                        f"Moved file to enqueued directory: {file_path.name} -> {unique_filename}"
                    )

                except Exception as move_error:
                    logger.error(
                        f"Failed to move file {file_path.name} to __enqueued__ directory: {move_error}"
                    )
                    # Don't affect the main function's success status

                return True, track_id

            except Exception as e:
                error_files = [
                    {
                        "file_path": str(file_path.name),
                        "error_description": "Document enqueue error",
                        "original_error": f"Failed to enqueue document: {str(e)}",
                        "file_size": file_size,
                    }
                ]
                await rag.apipeline_enqueue_error_documents(error_files, track_id)
                logger.error(f"Error enqueueing document {file_path.name}: {str(e)}")
                return False, track_id
        else:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "No content extracted",
                    "original_error": "No content could be extracted from file",
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(f"No content extracted from file: {file_path.name}")
            return False, track_id

    except Exception as e:
        # Catch-all for any unexpected errors
        try:
            file_size = file_path.stat().st_size if file_path.exists() else 0
        except Exception:
            file_size = 0

        error_files = [
            {
                "file_path": str(file_path.name),
                "error_description": "Unexpected processing error",
                "original_error": f"Unexpected error: {str(e)}",
                "file_size": file_size,
            }
        ]
        await rag.apipeline_enqueue_error_documents(error_files, track_id)
        logger.error(f"Enqueuing file {file_path.name} error: {str(e)}")
        logger.error(traceback.format_exc())
        return False, track_id
    finally:
        if file_path.name.startswith(temp_prefix):
            try:
                file_path.unlink()
            except Exception as e:
                logger.error(f"Error deleting file {file_path}: {str(e)}")


async def pipeline_index_file(rag: LightRAG, file_path: Path, track_id: str = None):
    """Index a file with track_id

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
        track_id: Optional tracking ID
    """
    try:
        success, returned_track_id = await pipeline_enqueue_file(
            rag, file_path, track_id
        )
        if success:
            await rag.apipeline_process_enqueue_documents()

    except Exception as e:
        logger.error(f"Error indexing file {file_path.name}: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_files(
    rag: LightRAG, file_paths: List[Path], track_id: str = None
):
    """Index multiple files sequentially to avoid high CPU load

    Args:
        rag: LightRAG instance
        file_paths: Paths to the files to index
        track_id: Optional tracking ID to pass to all files
    """
    if not file_paths:
        return
    try:
        enqueued = False

        # Use get_pinyin_sort_key for Chinese pinyin sorting
        sorted_file_paths = sorted(
            file_paths, key=lambda p: get_pinyin_sort_key(str(p))
        )

        # Process files sequentially with track_id
        for file_path in sorted_file_paths:
            success, _ = await pipeline_enqueue_file(rag, file_path, track_id)
            if success:
                enqueued = True

        # Process the queue only if at least one file was successfully enqueued
        if enqueued:
            await rag.apipeline_process_enqueue_documents()
    except Exception as e:
        logger.error(f"Error indexing files: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_texts(
    rag: LightRAG,
    texts: List[str],
    file_sources: List[str] = None,
    track_id: str = None,
):
    """Index a list of texts with track_id

    Args:
        rag: LightRAG instance
        texts: The texts to index
        file_sources: Sources of the texts
        track_id: Optional tracking ID
    """
    if not texts:
        return
    if file_sources is not None:
        if len(file_sources) != 0 and len(file_sources) != len(texts):
            [
                file_sources.append("unknown_source")
                for _ in range(len(file_sources), len(texts))
            ]
    await rag.apipeline_enqueue_documents(
        input=texts, file_paths=file_sources, track_id=track_id
    )
    await rag.apipeline_process_enqueue_documents()


async def run_scanning_process(
    rag: LightRAG, doc_manager: DocumentManager, track_id: str = None
):
    """Background task to scan and index documents

    Args:
        rag: LightRAG instance
        doc_manager: DocumentManager instance
        track_id: Optional tracking ID to pass to all scanned files
    """
    try:
        new_files = doc_manager.scan_directory_for_new_files()
        total_files = len(new_files)
        logger.info(f"Found {total_files} files to index.")

        if new_files:
            # Check for files with PROCESSED status and filter them out
            valid_files = []
            processed_files = []

            for file_path in new_files:
                filename = file_path.name
                existing_doc_data = await rag.doc_status.get_doc_by_file_path(filename)

                if existing_doc_data and existing_doc_data.get("status") == "processed":
                    # File is already PROCESSED, skip it with warning
                    processed_files.append(filename)
                    logger.warning(f"Skipping already processed file: {filename}")
                else:
                    # File is new or in non-PROCESSED status, add to processing list
                    valid_files.append(file_path)

            # Process valid files (new files + non-PROCESSED status files)
            if valid_files:
                await pipeline_index_files(rag, valid_files, track_id)
                if processed_files:
                    logger.info(
                        f"Scanning process completed: {len(valid_files)} files Processed {len(processed_files)} skipped."
                    )
                else:
                    logger.info(
                        f"Scanning process completed: {len(valid_files)} files Processed."
                    )
            else:
                logger.info(
                    "No files to process after filtering already processed files."
                )
        else:
            # No new files to index, check if there are any documents in the queue
            logger.info(
                "No upload file found, check if there are any documents in the queue..."
            )
            await rag.apipeline_process_enqueue_documents()

    except Exception as e:
        logger.error(f"Error during scanning process: {str(e)}")
        logger.error(traceback.format_exc())


async def background_delete_documents(
    rag: LightRAG,
    doc_manager: DocumentManager,
    doc_ids: List[str],
    delete_file: bool = False,
    delete_llm_cache: bool = False,
):
    """Background task to delete multiple documents"""
    from lightrag.kg.shared_storage import (
        get_namespace_data,
        get_namespace_lock,
    )

    pipeline_status = await get_namespace_data(
        "pipeline_status", workspace=rag.workspace
    )
    pipeline_status_lock = get_namespace_lock(
        "pipeline_status", workspace=rag.workspace
    )

    total_docs = len(doc_ids)
    successful_deletions = []
    failed_deletions = []

    # Double-check pipeline status before proceeding
    async with pipeline_status_lock:
        if pipeline_status.get("busy", False):
            logger.warning("Error: Unexpected pipeline busy state, aborting deletion.")
            return  # Abort deletion operation

        # Set pipeline status to busy for deletion
        pipeline_status.update(
            {
                "busy": True,
                # Job name can not be changed, it's verified in adelete_by_doc_id()
                "job_name": f"Deleting {total_docs} Documents",
                "job_start": datetime.now().isoformat(),
                "docs": total_docs,
                "batchs": total_docs,
                "cur_batch": 0,
                "latest_message": "Starting document deletion process",
            }
        )
        # Use slice assignment to clear the list in place
        pipeline_status["history_messages"][:] = ["Starting document deletion process"]
        if delete_llm_cache:
            pipeline_status["history_messages"].append(
                "LLM cache cleanup requested for this deletion job"
            )

    try:
        # Loop through each document ID and delete them one by one
        for i, doc_id in enumerate(doc_ids, 1):
            # Check for cancellation at the start of each document deletion
            async with pipeline_status_lock:
                if pipeline_status.get("cancellation_requested", False):
                    cancel_msg = f"Deletion cancelled by user at document {i}/{total_docs}. {len(successful_deletions)} deleted, {total_docs - i + 1} remaining."
                    logger.info(cancel_msg)
                    pipeline_status["latest_message"] = cancel_msg
                    pipeline_status["history_messages"].append(cancel_msg)
                    # Add remaining documents to failed list with cancellation reason
                    failed_deletions.extend(
                        doc_ids[i - 1 :]
                    )  # i-1 because enumerate starts at 1
                    break  # Exit the loop, remaining documents unchanged

                start_msg = f"Deleting document {i}/{total_docs}: {doc_id}"
                logger.info(start_msg)
                pipeline_status["cur_batch"] = i
                pipeline_status["latest_message"] = start_msg
                pipeline_status["history_messages"].append(start_msg)

            file_path = "#"
            try:
                result = await rag.adelete_by_doc_id(
                    doc_id, delete_llm_cache=delete_llm_cache
                )
                file_path = (
                    getattr(result, "file_path", "-") if "result" in locals() else "-"
                )
                if result.status == "success":
                    successful_deletions.append(doc_id)
                    success_msg = (
                        f"Document deleted {i}/{total_docs}: {doc_id}[{file_path}]"
                    )
                    logger.info(success_msg)
                    async with pipeline_status_lock:
                        pipeline_status["history_messages"].append(success_msg)

                    # Handle file deletion if requested and file_path is available
                    if (
                        delete_file
                        and result.file_path
                        and result.file_path != "unknown_source"
                    ):
                        try:
                            deleted_files = []
                            # SECURITY FIX: Use secure path validation to prevent arbitrary file deletion
                            safe_file_path = validate_file_path_security(
                                result.file_path, doc_manager.input_dir
                            )

                            if safe_file_path is None:
                                # Security violation detected - log and skip file deletion
                                security_msg = f"Security violation: Unsafe file path detected for deletion - {result.file_path}"
                                logger.warning(security_msg)
                                async with pipeline_status_lock:
                                    pipeline_status["latest_message"] = security_msg
                                    pipeline_status["history_messages"].append(
                                        security_msg
                                    )
                            else:
                                # check and delete files from input_dir directory
                                if safe_file_path.exists():
                                    try:
                                        safe_file_path.unlink()
                                        deleted_files.append(safe_file_path.name)
                                        file_delete_msg = f"Successfully deleted input_dir file: {result.file_path}"
                                        logger.info(file_delete_msg)
                                        async with pipeline_status_lock:
                                            pipeline_status["latest_message"] = (
                                                file_delete_msg
                                            )
                                            pipeline_status["history_messages"].append(
                                                file_delete_msg
                                            )
                                    except Exception as file_error:
                                        file_error_msg = f"Failed to delete input_dir file {result.file_path}: {str(file_error)}"
                                        logger.debug(file_error_msg)
                                        async with pipeline_status_lock:
                                            pipeline_status["latest_message"] = (
                                                file_error_msg
                                            )
                                            pipeline_status["history_messages"].append(
                                                file_error_msg
                                            )

                                # Also check and delete files from __enqueued__ directory
                                enqueued_dir = doc_manager.input_dir / "__enqueued__"
                                if enqueued_dir.exists():
                                    # SECURITY FIX: Validate that the file path is safe before processing
                                    # Only proceed if the original path validation passed
                                    base_name = Path(result.file_path).stem
                                    extension = Path(result.file_path).suffix

                                    # Search for exact match and files with numeric suffixes
                                    for enqueued_file in enqueued_dir.glob(
                                        f"{base_name}*{extension}"
                                    ):
                                        # Additional security check: ensure enqueued file is within enqueued directory
                                        safe_enqueued_path = (
                                            validate_file_path_security(
                                                enqueued_file.name, enqueued_dir
                                            )
                                        )
                                        if safe_enqueued_path is not None:
                                            try:
                                                enqueued_file.unlink()
                                                deleted_files.append(enqueued_file.name)
                                                logger.info(
                                                    f"Successfully deleted enqueued file: {enqueued_file.name}"
                                                )
                                            except Exception as enqueued_error:
                                                file_error_msg = f"Failed to delete enqueued file {enqueued_file.name}: {str(enqueued_error)}"
                                                logger.debug(file_error_msg)
                                                async with pipeline_status_lock:
                                                    pipeline_status[
                                                        "latest_message"
                                                    ] = file_error_msg
                                                    pipeline_status[
                                                        "history_messages"
                                                    ].append(file_error_msg)
                                        else:
                                            security_msg = f"Security violation: Unsafe enqueued file path detected - {enqueued_file.name}"
                                            logger.warning(security_msg)

                            if deleted_files == []:
                                file_error_msg = f"File deletion skipped, missing or unsafe file: {result.file_path}"
                                logger.warning(file_error_msg)
                                async with pipeline_status_lock:
                                    pipeline_status["latest_message"] = file_error_msg
                                    pipeline_status["history_messages"].append(
                                        file_error_msg
                                    )

                        except Exception as file_error:
                            file_error_msg = f"Failed to delete file {result.file_path}: {str(file_error)}"
                            logger.error(file_error_msg)
                            async with pipeline_status_lock:
                                pipeline_status["latest_message"] = file_error_msg
                                pipeline_status["history_messages"].append(
                                    file_error_msg
                                )
                    elif delete_file:
                        no_file_msg = (
                            f"File deletion skipped, missing file path: {doc_id}"
                        )
                        logger.warning(no_file_msg)
                        async with pipeline_status_lock:
                            pipeline_status["latest_message"] = no_file_msg
                            pipeline_status["history_messages"].append(no_file_msg)
                else:
                    failed_deletions.append(doc_id)
                    error_msg = f"Failed to delete {i}/{total_docs}: {doc_id}[{file_path}] - {result.message}"
                    logger.error(error_msg)
                    async with pipeline_status_lock:
                        pipeline_status["latest_message"] = error_msg
                        pipeline_status["history_messages"].append(error_msg)

            except Exception as e:
                failed_deletions.append(doc_id)
                error_msg = f"Error deleting document {i}/{total_docs}: {doc_id}[{file_path}] - {str(e)}"
                logger.error(error_msg)
                logger.error(traceback.format_exc())
                async with pipeline_status_lock:
                    pipeline_status["latest_message"] = error_msg
                    pipeline_status["history_messages"].append(error_msg)

    except Exception as e:
        error_msg = f"Critical error during batch deletion: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        async with pipeline_status_lock:
            pipeline_status["history_messages"].append(error_msg)
    finally:
        # Final summary and check for pending requests
        async with pipeline_status_lock:
            pipeline_status["busy"] = False
            pipeline_status["pending_requests"] = False  # Reset pending requests flag
            pipeline_status["cancellation_requested"] = (
                False  # Always reset cancellation flag
            )
            completion_msg = f"Deletion completed: {len(successful_deletions)} successful, {len(failed_deletions)} failed"
            pipeline_status["latest_message"] = completion_msg
            pipeline_status["history_messages"].append(completion_msg)

            # Check if there are pending document indexing requests
            has_pending_request = pipeline_status.get("request_pending", False)

        # If there are pending requests, start document processing pipeline
        if has_pending_request:
            try:
                logger.info(
                    "Processing pending document indexing requests after deletion"
                )
                await rag.apipeline_process_enqueue_documents()
            except Exception as e:
                logger.error(f"Error processing pending documents after deletion: {e}")


def create_document_routes(
    rag: LightRAG, doc_manager: DocumentManager, api_key: Optional[str] = None
):
    # Create combined auth dependency for document routes
    combined_auth = get_combined_auth_dependency(api_key)

    @router.post(
        "/scan", response_model=ScanResponse, dependencies=[Depends(combined_auth)],
        summary="扫描新文档",
        description="触发扫描新文档的进程。",
        response_description="扫描操作的结果"
    )
    async def scan_for_new_documents(background_tasks: BackgroundTasks):
        """
        触发新文档的扫描进程。

        此端点启动一个后台任务，扫描输入目录中的新文档并处理它们。
        如果扫描进程已在运行，则返回指示该事实的状态。

        返回:
            ScanResponse: 包含扫描状态和track_id的响应对象
        """
        # 生成带"scan"前缀的track_id用于扫描操作
        track_id = generate_track_id("scan")

        # 使用track_id在后台启动扫描进程
        background_tasks.add_task(run_scanning_process, rag, doc_manager, track_id)
        return ScanResponse(
            status="scanning_started",
            message="后台扫描进程已启动",
            track_id=track_id,
        )

    @router.post(
        "/upload", response_model=InsertResponse, dependencies=[Depends(combined_auth)],
        summary="上传文档",
        description="上传文件到输入目录并建立索引。",
        response_description="上传操作的结果"
    )
    async def upload_to_input_dir(
        background_tasks: BackgroundTasks, file: UploadFile = File(...)
    ):
        """
        上传文件到输入目录并建立索引。

        此API端点通过HTTP POST请求接收文件，检查上传的文件是否为支持的类型，
        将其保存在指定的输入目录中，建立索引以供检索，并返回包含相关详细信息的成功状态。

        参数:
            background_tasks: FastAPI后台任务用于异步处理
            file (UploadFile): 要上传的文件。它必须具有允许的扩展名。

        返回:
            InsertResponse: 包含上传状态和消息的响应对象。
                status可以是"success"、"duplicated"，或者抛出错误。

        异常:
            HTTPException: 如果文件类型不受支持（400）或其他错误发生（500）。
        """
        try:
            # Sanitize filename to prevent Path Traversal attacks
            safe_filename = sanitize_filename(file.filename, doc_manager.input_dir)

            if not doc_manager.is_supported_file(safe_filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            # Check if filename already exists in doc_status storage
            existing_doc_data = await rag.doc_status.get_doc_by_file_path(safe_filename)
            if existing_doc_data:
                # Get document status and track_id from existing document
                status = existing_doc_data.get("status", "unknown")
                # Use `or ""` to handle both missing key and None value (e.g., legacy rows without track_id)
                existing_track_id = existing_doc_data.get("track_id") or ""
                return InsertResponse(
                    status="duplicated",
                    message=f"File '{safe_filename}' already exists in document storage (Status: {status}).",
                    track_id=existing_track_id,
                )

            file_path = doc_manager.input_dir / safe_filename
            # Check if file already exists in file system
            if file_path.exists():
                return InsertResponse(
                    status="duplicated",
                    message=f"File '{safe_filename}' already exists in the input directory.",
                    track_id="",
                )

            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            track_id = generate_track_id("upload")

            # Add to background tasks and get track_id
            background_tasks.add_task(pipeline_index_file, rag, file_path, track_id)

            return InsertResponse(
                status="success",
                message=f"File '{safe_filename}' uploaded successfully. Processing will continue in background.",
                track_id=track_id,
            )

        except Exception as e:
            logger.error(f"Error /documents/upload: {file.filename}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/text", response_model=InsertResponse, dependencies=[Depends(combined_auth)],
        summary="插入文本",
        description="将文本插入到RAG系统中。",
        response_description="文本插入操作的结果"
    )
    async def insert_text(
        request: InsertTextRequest, background_tasks: BackgroundTasks
    ):
        """
        将文本插入到RAG系统中。

        此端点允许您将文本数据插入到RAG系统中，以供后续检索并在生成响应时使用。

        参数:
            request (InsertTextRequest): 包含要插入文本的请求体。
            background_tasks: FastAPI后台任务用于异步处理

        返回:
            InsertResponse: 包含操作状态的响应对象。

        异常:
            HTTPException: 如果在文本处理过程中发生错误（500）。
        """
        try:
            # Check if file_source already exists in doc_status storage
            if (
                request.file_source
                and request.file_source.strip()
                and request.file_source != "unknown_source"
            ):
                existing_doc_data = await rag.doc_status.get_doc_by_file_path(
                    request.file_source
                )
                if existing_doc_data:
                    # Get document status and track_id from existing document
                    status = existing_doc_data.get("status", "unknown")
                    # Use `or ""` to handle both missing key and None value (e.g., legacy rows without track_id)
                    existing_track_id = existing_doc_data.get("track_id") or ""
                    return InsertResponse(
                        status="duplicated",
                        message=f"File source '{request.file_source}' already exists in document storage (Status: {status}).",
                        track_id=existing_track_id,
                    )

            # Check if content already exists by computing content hash (doc_id)
            sanitized_text = sanitize_text_for_encoding(request.text)
            content_doc_id = compute_mdhash_id(sanitized_text, prefix="doc-")
            existing_doc = await rag.doc_status.get_by_id(content_doc_id)
            if existing_doc:
                # Content already exists, return duplicated with existing track_id
                status = existing_doc.get("status", "unknown")
                existing_track_id = existing_doc.get("track_id") or ""
                return InsertResponse(
                    status="duplicated",
                    message=f"Identical content already exists in document storage (doc_id: {content_doc_id}, Status: {status}).",
                    track_id=existing_track_id,
                )

            # Generate track_id for text insertion
            track_id = generate_track_id("insert")

            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                [request.text],
                file_sources=[request.file_source],
                track_id=track_id,
            )

            return InsertResponse(
                status="success",
                message="Text successfully received. Processing will continue in background.",
                track_id=track_id,
            )
        except Exception as e:
            logger.error(f"Error /documents/text: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/texts",
        response_model=InsertResponse,
        dependencies=[Depends(combined_auth)],
        summary="插入多个文本",
        description="将多个文本插入到RAG系统中。",
        response_description="文本插入操作的结果"
    )
    async def insert_texts(
        request: InsertTextsRequest, background_tasks: BackgroundTasks
    ):
        """
        将多个文本插入到RAG系统中。

        此端点允许您将多个文本数据插入到RAG系统中，以供后续检索并在生成响应时使用。

        参数:
            request (InsertTextsRequest): 包含要插入文本的请求体。
            background_tasks: FastAPI后台任务用于异步处理

        返回:
            InsertResponse: 包含操作状态的响应对象。

        异常:
            HTTPException: 如果在文本处理过程中发生错误（500）。
        """
        try:
            # 检查file_sources是否已存在于doc_status存储中
            if request.file_sources:
                for file_source in request.file_sources:
                    if file_source and file_source.strip() and file_source != "unknown_source":
                        existing_doc_data = await rag.doc_status.get_doc_by_file_path(
                            file_source
                        )
                        if existing_doc_data:
                            # 从现有文档获取文档状态和track_id
                            status = existing_doc_data.get("status", "unknown")
                            # 使用`or ""`处理缺失键和None值（例如，没有track_id的旧行）
                            existing_track_id = existing_doc_data.get("track_id") or ""
                            return InsertResponse(
                                status="duplicated",
                                message=f"文本源'{file_source}'已存在于文档存储中（状态: {status}）。",
                                track_id=existing_track_id,
                            )

            # 通过计算内容哈希（doc_id）检查任何内容是否已存在
            for text in request.texts:
                sanitized_text = sanitize_text_for_encoding(text)
                content_doc_id = compute_mdhash_id(sanitized_text, prefix="doc-")
                existing_doc = await rag.doc_status.get_by_id(content_doc_id)
                if existing_doc:
                    # 内容已存在，返回重复并附带现有track_id
                    status = existing_doc.get("status", "unknown")
                    existing_track_id = existing_doc.get("track_id") or ""
                    return InsertResponse(
                        status="duplicated",
                        message=f"相同内容已存在于文档存储中（doc_id: {content_doc_id}，状态: {status}）。",
                        track_id=existing_track_id,
                    )

            # 生成跟踪ID用于文本插入
            track_id = generate_track_id("insert")

            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                request.texts,
                file_sources=request.file_sources,
                track_id=track_id,
            )

            return InsertResponse(
                status="success",
                message=f"{len(request.texts)}个文本已成功提交处理。将在后台继续处理。",
                track_id=track_id,
            )
        except Exception as e:
            logger.error(f"插入文本时出错: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.delete(
        "", response_model=ClearDocumentsResponse, dependencies=[Depends(combined_auth)]
    )
    async def clear_documents():
        """
        Clear all documents from the RAG system.

        This endpoint deletes all documents, entities, relationships, and files from the system.
        It uses the storage drop methods to properly clean up all data and removes all files
        from the input directory.

        Returns:
            ClearDocumentsResponse: A response object containing the status and message.
                - status="success":           All documents and files were successfully cleared.
                - status="partial_success":   Document clear job exit with some errors.
                - status="busy":              Operation could not be completed because the pipeline is busy.
                - status="fail":              All storage drop operations failed, with message
                - message: Detailed information about the operation results, including counts
                  of deleted files and any errors encountered.

        Raises:
            HTTPException: Raised when a serious error occurs during the clearing process,
                          with status code 500 and error details in the detail field.
        """
        from lightrag.kg.shared_storage import (
            get_namespace_data,
            get_namespace_lock,
        )

        # Get pipeline status and lock
        pipeline_status = await get_namespace_data(
            "pipeline_status", workspace=rag.workspace
        )
        pipeline_status_lock = get_namespace_lock(
            "pipeline_status", workspace=rag.workspace
        )

        # Check and set status with lock
        async with pipeline_status_lock:
            if pipeline_status.get("busy", False):
                return ClearDocumentsResponse(
                    status="busy",
                    message="Cannot clear documents while pipeline is busy",
                )
            # Set busy to true
            pipeline_status.update(
                {
                    "busy": True,
                    "job_name": "Clearing Documents",
                    "job_start": datetime.now().isoformat(),
                    "docs": 0,
                    "batchs": 0,
                    "cur_batch": 0,
                    "request_pending": False,  # Clear any previous request
                    "latest_message": "Starting document clearing process",
                }
            )
            # Cleaning history_messages without breaking it as a shared list object
            del pipeline_status["history_messages"][:]
            pipeline_status["history_messages"].append(
                "Starting document clearing process"
            )

        try:
            # Use drop method to clear all data
            drop_tasks = []
            storages = [
                rag.text_chunks,
                rag.full_docs,
                rag.full_entities,
                rag.full_relations,
                rag.entity_chunks,
                rag.relation_chunks,
                rag.entities_vdb,
                rag.relationships_vdb,
                rag.chunks_vdb,
                rag.chunk_entity_relation_graph,
                rag.doc_status,
            ]

            # Log storage drop start
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(
                    "Starting to drop storage components"
                )

            for storage in storages:
                if storage is not None:
                    drop_tasks.append(storage.drop())

            # Wait for all drop tasks to complete
            drop_results = await asyncio.gather(*drop_tasks, return_exceptions=True)

            # Check for errors and log results
            errors = []
            storage_success_count = 0
            storage_error_count = 0

            for i, result in enumerate(drop_results):
                storage_name = storages[i].__class__.__name__
                if isinstance(result, Exception):
                    error_msg = f"Error dropping {storage_name}: {str(result)}"
                    errors.append(error_msg)
                    logger.error(error_msg)
                    storage_error_count += 1
                else:
                    namespace = storages[i].namespace
                    workspace = storages[i].workspace
                    logger.info(
                        f"Successfully dropped {storage_name}: {workspace}/{namespace}"
                    )
                    storage_success_count += 1

            # Log storage drop results
            if "history_messages" in pipeline_status:
                if storage_error_count > 0:
                    pipeline_status["history_messages"].append(
                        f"Dropped {storage_success_count} storage components with {storage_error_count} errors"
                    )
                else:
                    pipeline_status["history_messages"].append(
                        f"Successfully dropped all {storage_success_count} storage components"
                    )

            # If all storage operations failed, return error status and don't proceed with file deletion
            if storage_success_count == 0 and storage_error_count > 0:
                error_message = "All storage drop operations failed. Aborting document clearing process."
                logger.error(error_message)
                if "history_messages" in pipeline_status:
                    pipeline_status["history_messages"].append(error_message)
                return ClearDocumentsResponse(status="fail", message=error_message)

            # Log file deletion start
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(
                    "Starting to delete files in input directory"
                )

            # Delete only files in the current directory, preserve files in subdirectories
            deleted_files_count = 0
            file_errors_count = 0

            for file_path in doc_manager.input_dir.glob("*"):
                if file_path.is_file():
                    try:
                        file_path.unlink()
                        deleted_files_count += 1
                    except Exception as e:
                        logger.error(f"Error deleting file {file_path}: {str(e)}")
                        file_errors_count += 1

            # Log file deletion results
            if "history_messages" in pipeline_status:
                if file_errors_count > 0:
                    pipeline_status["history_messages"].append(
                        f"Deleted {deleted_files_count} files with {file_errors_count} errors"
                    )
                    errors.append(f"Failed to delete {file_errors_count} files")
                else:
                    pipeline_status["history_messages"].append(
                        f"Successfully deleted {deleted_files_count} files"
                    )

            # Prepare final result message
            final_message = ""
            if errors:
                final_message = f"Cleared documents with some errors. Deleted {deleted_files_count} files."
                status = "partial_success"
            else:
                final_message = f"All documents cleared successfully. Deleted {deleted_files_count} files."
                status = "success"

            # Log final result
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(final_message)

            # Return response based on results
            return ClearDocumentsResponse(status=status, message=final_message)
        except Exception as e:
            error_msg = f"Error clearing documents: {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(error_msg)
            raise HTTPException(status_code=500, detail=str(e))
        finally:
            # Reset busy status after completion
            async with pipeline_status_lock:
                pipeline_status["busy"] = False
                completion_msg = "Document clearing process completed"
                pipeline_status["latest_message"] = completion_msg
                if "history_messages" in pipeline_status:
                    pipeline_status["history_messages"].append(completion_msg)

    @router.get(
        "/pipeline_status",
        dependencies=[Depends(combined_auth)],
        response_model=PipelineStatusResponse,
    )
    async def get_pipeline_status() -> PipelineStatusResponse:
        """
        Get the current status of the document indexing pipeline.

        This endpoint returns information about the current state of the document processing pipeline,
        including the processing status, progress information, and history messages.

        Returns:
            PipelineStatusResponse: A response object containing:
                - autoscanned (bool): Whether auto-scan has started
                - busy (bool): Whether the pipeline is currently busy
                - job_name (str): Current job name (e.g., indexing files/indexing texts)
                - job_start (str, optional): Job start time as ISO format string
                - docs (int): Total number of documents to be indexed
                - batchs (int): Number of batches for processing documents
                - cur_batch (int): Current processing batch
                - request_pending (bool): Flag for pending request for processing
                - latest_message (str): Latest message from pipeline processing
                - history_messages (List[str], optional): List of history messages (limited to latest 1000 entries,
                  with truncation message if more than 1000 messages exist)

        Raises:
            HTTPException: If an error occurs while retrieving pipeline status (500)
        """
        try:
            from lightrag.kg.shared_storage import (
                get_namespace_data,
                get_namespace_lock,
                get_all_update_flags_status,
            )

            pipeline_status = await get_namespace_data(
                "pipeline_status", workspace=rag.workspace
            )
            pipeline_status_lock = get_namespace_lock(
                "pipeline_status", workspace=rag.workspace
            )

            # Get update flags status for all namespaces
            update_status = await get_all_update_flags_status(workspace=rag.workspace)

            # Convert MutableBoolean objects to regular boolean values
            processed_update_status = {}
            for namespace, flags in update_status.items():
                processed_flags = []
                for flag in flags:
                    # Handle both multiprocess and single process cases
                    if hasattr(flag, "value"):
                        processed_flags.append(bool(flag.value))
                    else:
                        processed_flags.append(bool(flag))
                processed_update_status[namespace] = processed_flags

            async with pipeline_status_lock:
                # Convert to regular dict if it's a Manager.dict
                status_dict = dict(pipeline_status)

            # Add processed update_status to the status dictionary
            status_dict["update_status"] = processed_update_status

            # Convert history_messages to a regular list if it's a Manager.list
            # and limit to latest 1000 entries with truncation message if needed
            if "history_messages" in status_dict:
                history_list = list(status_dict["history_messages"])
                total_count = len(history_list)

                if total_count > 1000:
                    # Calculate truncated message count
                    truncated_count = total_count - 1000

                    # Take only the latest 1000 messages
                    latest_messages = history_list[-1000:]

                    # Add truncation message at the beginning
                    truncation_message = (
                        f"[Truncated history messages: {truncated_count}/{total_count}]"
                    )
                    status_dict["history_messages"] = [
                        truncation_message
                    ] + latest_messages
                else:
                    # No truncation needed, return all messages
                    status_dict["history_messages"] = history_list

            # Ensure job_start is properly formatted as a string with timezone information
            if "job_start" in status_dict and status_dict["job_start"]:
                # Use format_datetime to ensure consistent formatting
                status_dict["job_start"] = format_datetime(status_dict["job_start"])

            return PipelineStatusResponse(**status_dict)
        except Exception as e:
            logger.error(f"Error getting pipeline status: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    # TODO: Deprecated, use /documents/paginated instead
    @router.get(
        "", response_model=DocsStatusesResponse, dependencies=[Depends(combined_auth)]
    )
    async def documents() -> DocsStatusesResponse:
        """
        Get the status of all documents in the system. This endpoint is deprecated; use /documents/paginated instead.
        To prevent excessive resource consumption, a maximum of 1,000 records is returned.

        This endpoint retrieves the current status of all documents, grouped by their
        processing status (PENDING, PROCESSING, PREPROCESSED, PROCESSED, FAILED). The results are
        limited to 1000 total documents with fair distribution across all statuses.

        Returns:
            DocsStatusesResponse: A response object containing a dictionary where keys are
                                DocStatus values and values are lists of DocStatusResponse
                                objects representing documents in each status category.
                                Maximum 1000 documents total will be returned.

        Raises:
            HTTPException: If an error occurs while retrieving document statuses (500).
        """
        try:
            statuses = (
                DocStatus.PENDING,
                DocStatus.PROCESSING,
                DocStatus.PREPROCESSED,
                DocStatus.PROCESSED,
                DocStatus.FAILED,
            )

            tasks = [rag.get_docs_by_status(status) for status in statuses]
            results: List[Dict[str, DocProcessingStatus]] = await asyncio.gather(*tasks)

            response = DocsStatusesResponse()
            total_documents = 0
            max_documents = 1000

            # Convert results to lists for easier processing
            status_documents = []
            for idx, result in enumerate(results):
                status = statuses[idx]
                docs_list = []
                for doc_id, doc_status in result.items():
                    docs_list.append((doc_id, doc_status))
                status_documents.append((status, docs_list))

            # Fair distribution: round-robin across statuses
            status_indices = [0] * len(
                status_documents
            )  # Track current index for each status
            current_status_idx = 0

            while total_documents < max_documents:
                # Check if we have any documents left to process
                has_remaining = False
                for status_idx, (status, docs_list) in enumerate(status_documents):
                    if status_indices[status_idx] < len(docs_list):
                        has_remaining = True
                        break

                if not has_remaining:
                    break

                # Try to get a document from the current status
                status, docs_list = status_documents[current_status_idx]
                current_index = status_indices[current_status_idx]

                if current_index < len(docs_list):
                    doc_id, doc_status = docs_list[current_index]

                    if status not in response.statuses:
                        response.statuses[status] = []

                    response.statuses[status].append(
                        DocStatusResponse(
                            id=doc_id,
                            content_summary=doc_status.content_summary,
                            content_length=doc_status.content_length,
                            status=doc_status.status,
                            created_at=format_datetime(doc_status.created_at),
                            updated_at=format_datetime(doc_status.updated_at),
                            track_id=doc_status.track_id,
                            chunks_count=doc_status.chunks_count,
                            error_msg=doc_status.error_msg,
                            metadata=doc_status.metadata,
                            file_path=doc_status.file_path,
                        )
                    )

                    status_indices[current_status_idx] += 1
                    total_documents += 1

                # Move to next status (round-robin)
                current_status_idx = (current_status_idx + 1) % len(status_documents)

            return response
        except Exception as e:
            logger.error(f"Error GET /documents: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    class DeleteDocByIdResponse(BaseModel):
        """Response model for single document deletion operation."""

        status: Literal["deletion_started", "busy", "not_allowed"] = Field(
            description="Status of the deletion operation"
        )
        message: str = Field(description="Message describing the operation result")
        doc_id: str = Field(description="The ID of the document to delete")

    @router.delete(
        "/delete_document",
        response_model=DeleteDocByIdResponse,
        dependencies=[Depends(combined_auth)],
        summary="Delete a document and all its associated data by its ID.",
    )
    async def delete_document(
        delete_request: DeleteDocRequest,
        background_tasks: BackgroundTasks,
    ) -> DeleteDocByIdResponse:
        """
        Delete documents and all their associated data by their IDs using background processing.

        Deletes specific documents and all their associated data, including their status,
        text chunks, vector embeddings, and any related graph data. When requested,
        cached LLM extraction responses are removed after graph deletion/rebuild completes.
        The deletion process runs in the background to avoid blocking the client connection.

        This operation is irreversible and will interact with the pipeline status.

        Args:
            delete_request (DeleteDocRequest): The request containing the document IDs and deletion options.
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            DeleteDocByIdResponse: The result of the deletion operation.
                - status="deletion_started": The document deletion has been initiated in the background.
                - status="busy": The pipeline is busy with another operation.

        Raises:
            HTTPException:
              - 500: If an unexpected internal error occurs during initialization.
        """
        doc_ids = delete_request.doc_ids

        try:
            from lightrag.kg.shared_storage import (
                get_namespace_data,
                get_namespace_lock,
            )

            pipeline_status = await get_namespace_data(
                "pipeline_status", workspace=rag.workspace
            )
            pipeline_status_lock = get_namespace_lock(
                "pipeline_status", workspace=rag.workspace
            )

            # Check if pipeline is busy with proper lock
            async with pipeline_status_lock:
                if pipeline_status.get("busy", False):
                    return DeleteDocByIdResponse(
                        status="busy",
                        message="Cannot delete documents while pipeline is busy",
                        doc_id=", ".join(doc_ids),
                    )

            # Add deletion task to background tasks
            background_tasks.add_task(
                background_delete_documents,
                rag,
                doc_manager,
                doc_ids,
                delete_request.delete_file,
                delete_request.delete_llm_cache,
            )

            return DeleteDocByIdResponse(
                status="deletion_started",
                message=f"Document deletion for {len(doc_ids)} documents has been initiated. Processing will continue in background.",
                doc_id=", ".join(doc_ids),
            )

        except Exception as e:
            error_msg = f"Error initiating document deletion for {delete_request.doc_ids}: {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=error_msg)

    @router.post(
        "/clear_cache",
        response_model=ClearCacheResponse,
        dependencies=[Depends(combined_auth)],
        summary="清除缓存",
        description="清除LLM响应缓存存储中的所有缓存数据。",
        response_description="清除缓存操作的结果"
    )
    async def clear_cache(request: ClearCacheRequest):
        """
        清除LLM响应缓存存储中的所有缓存数据。

        此端点清除所有模式的缓存LLM响应。
        出于API兼容性考虑接受请求体，但会被忽略。

        参数:
            request (ClearCacheRequest): 请求体（出于兼容性考虑被忽略）。

        返回:
            ClearCacheResponse: 包含状态和消息的响应对象。

        异常:
            HTTPException: 如果清除缓存时发生错误（500）。
        """
        try:
            # 调用aclear_cache方法（无模式参数）
            await rag.aclear_cache()

            # 准备成功消息
            message = "已成功清除所有缓存"

            return ClearCacheResponse(status="success", message=message)
        except Exception as e:
            logger.error(f"清除缓存时出错: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.delete(
        "/delete_entity",
        response_model=DeletionResult,
        dependencies=[Depends(combined_auth)],
    )
    async def delete_entity(request: DeleteEntityRequest):
        """
        Delete an entity and all its relationships from the knowledge graph.

        Args:
            request (DeleteEntityRequest): The request body containing the entity name.

        Returns:
            DeletionResult: An object containing the outcome of the deletion process.

        Raises:
            HTTPException: If the entity is not found (404) or an error occurs (500).
        """
        try:
            result = await rag.adelete_by_entity(entity_name=request.entity_name)
            if result.status == "not_found":
                raise HTTPException(status_code=404, detail=result.message)
            if result.status == "fail":
                raise HTTPException(status_code=500, detail=result.message)
            # Set doc_id to empty string since this is an entity operation, not document
            result.doc_id = ""
            return result
        except HTTPException:
            raise
        except Exception as e:
            error_msg = f"Error deleting entity '{request.entity_name}': {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=error_msg)

    @router.delete(
        "/delete_relation",
        response_model=DeletionResult,
        dependencies=[Depends(combined_auth)],
    )
    async def delete_relation(request: DeleteRelationRequest):
        """
        Delete a relationship between two entities from the knowledge graph.

        Args:
            request (DeleteRelationRequest): The request body containing the source and target entity names.

        Returns:
            DeletionResult: An object containing the outcome of the deletion process.

        Raises:
            HTTPException: If the relation is not found (404) or an error occurs (500).
        """
        try:
            result = await rag.adelete_by_relation(
                source_entity=request.source_entity,
                target_entity=request.target_entity,
            )
            if result.status == "not_found":
                raise HTTPException(status_code=404, detail=result.message)
            if result.status == "fail":
                raise HTTPException(status_code=500, detail=result.message)
            # Set doc_id to empty string since this is a relation operation, not document
            result.doc_id = ""
            return result
        except HTTPException:
            raise
        except Exception as e:
            error_msg = f"Error deleting relation from '{request.source_entity}' to '{request.target_entity}': {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=error_msg)

    @router.get(
        "/track_status/{track_id}",
        response_model=TrackStatusResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def get_track_status(track_id: str) -> TrackStatusResponse:
        """
        Get the processing status of documents by tracking ID.

        This endpoint retrieves all documents associated with a specific tracking ID,
        allowing users to monitor the processing progress of their uploaded files or inserted texts.

        Args:
            track_id (str): The tracking ID returned from upload, text, or texts endpoints

        Returns:
            TrackStatusResponse: A response object containing:
                - track_id: The tracking ID
                - documents: List of documents associated with this track_id
                - total_count: Total number of documents for this track_id

        Raises:
            HTTPException: If track_id is invalid (400) or an error occurs (500).
        """
        try:
            # Validate track_id
            if not track_id or not track_id.strip():
                raise HTTPException(status_code=400, detail="Track ID cannot be empty")

            track_id = track_id.strip()

            # Get documents by track_id
            docs_by_track_id = await rag.aget_docs_by_track_id(track_id)

            # Convert to response format
            documents = []
            status_summary = {}

            for doc_id, doc_status in docs_by_track_id.items():
                documents.append(
                    DocStatusResponse(
                        id=doc_id,
                        content_summary=doc_status.content_summary,
                        content_length=doc_status.content_length,
                        status=doc_status.status,
                        created_at=format_datetime(doc_status.created_at),
                        updated_at=format_datetime(doc_status.updated_at),
                        track_id=doc_status.track_id,
                        chunks_count=doc_status.chunks_count,
                        error_msg=doc_status.error_msg,
                        metadata=doc_status.metadata,
                        file_path=doc_status.file_path,
                    )
                )

                # Build status summary
                # Handle both DocStatus enum and string cases for robust deserialization
                status_key = str(doc_status.status)
                status_summary[status_key] = status_summary.get(status_key, 0) + 1

            return TrackStatusResponse(
                track_id=track_id,
                documents=documents,
                total_count=len(documents),
                status_summary=status_summary,
            )

        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error getting track status for {track_id}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/paginated",
        response_model=PaginatedDocsResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def get_documents_paginated(
        request: DocumentsRequest,
    ) -> PaginatedDocsResponse:
        """
        Get documents with pagination support.

        This endpoint retrieves documents with pagination, filtering, and sorting capabilities.
        It provides better performance for large document collections by loading only the
        requested page of data.

        Args:
            request (DocumentsRequest): The request body containing pagination parameters

        Returns:
            PaginatedDocsResponse: A response object containing:
                - documents: List of documents for the current page
                - pagination: Pagination information (page, total_count, etc.)
                - status_counts: Count of documents by status for all documents

        Raises:
            HTTPException: If an error occurs while retrieving documents (500).
        """
        try:
            # Get paginated documents and status counts in parallel
            docs_task = rag.doc_status.get_docs_paginated(
                status_filter=request.status_filter,
                page=request.page,
                page_size=request.page_size,
                sort_field=request.sort_field,
                sort_direction=request.sort_direction,
            )
            status_counts_task = rag.doc_status.get_all_status_counts()

            # Execute both queries in parallel
            (documents_with_ids, total_count), status_counts = await asyncio.gather(
                docs_task, status_counts_task
            )

            # Convert documents to response format
            doc_responses = []
            for doc_id, doc in documents_with_ids:
                doc_responses.append(
                    DocStatusResponse(
                        id=doc_id,
                        content_summary=doc.content_summary,
                        content_length=doc.content_length,
                        status=doc.status,
                        created_at=format_datetime(doc.created_at),
                        updated_at=format_datetime(doc.updated_at),
                        track_id=doc.track_id,
                        chunks_count=doc.chunks_count,
                        error_msg=doc.error_msg,
                        metadata=doc.metadata,
                        file_path=doc.file_path,
                    )
                )

            # Calculate pagination info
            total_pages = (total_count + request.page_size - 1) // request.page_size
            has_next = request.page < total_pages
            has_prev = request.page > 1

            pagination = PaginationInfo(
                page=request.page,
                page_size=request.page_size,
                total_count=total_count,
                total_pages=total_pages,
                has_next=has_next,
                has_prev=has_prev,
            )

            return PaginatedDocsResponse(
                documents=doc_responses,
                pagination=pagination,
                status_counts=status_counts,
            )

        except Exception as e:
            logger.error(f"Error getting paginated documents: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.get(
        "/status_counts",
        response_model=StatusCountsResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def get_document_status_counts() -> StatusCountsResponse:
        """
        Get counts of documents by status.

        This endpoint retrieves the count of documents in each processing status
        (PENDING, PROCESSING, PROCESSED, FAILED) for all documents in the system.

        Returns:
            StatusCountsResponse: A response object containing status counts

        Raises:
            HTTPException: If an error occurs while retrieving status counts (500).
        """
        try:
            status_counts = await rag.doc_status.get_all_status_counts()
            return StatusCountsResponse(status_counts=status_counts)

        except Exception as e:
            logger.error(f"Error getting document status counts: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/reprocess_failed",
        response_model=ReprocessResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def reprocess_failed_documents(background_tasks: BackgroundTasks):
        """
        Reprocess failed and pending documents.

        This endpoint triggers the document processing pipeline which automatically
        picks up and reprocesses documents in the following statuses:
        - FAILED: Documents that failed during previous processing attempts
        - PENDING: Documents waiting to be processed
        - PROCESSING: Documents with abnormally terminated processing (e.g., server crashes)

        This is useful for recovering from server crashes, network errors, LLM service
        outages, or other temporary failures that caused document processing to fail.

        The processing happens in the background and can be monitored by checking the
        pipeline status. The reprocessed documents retain their original track_id from
        initial upload, so use their original track_id to monitor progress.

        Returns:
            ReprocessResponse: Response with status and message.
                track_id is always empty string because reprocessed documents retain
                their original track_id from initial upload.

        Raises:
            HTTPException: If an error occurs while initiating reprocessing (500).
        """
        try:
            # Start the reprocessing in the background
            # Note: Reprocessed documents retain their original track_id from initial upload
            background_tasks.add_task(rag.apipeline_process_enqueue_documents)
            logger.info("Reprocessing of failed documents initiated")

            return ReprocessResponse(
                status="reprocessing_started",
                message="Reprocessing of failed documents has been initiated in background. Documents retain their original track_id.",
            )

        except Exception as e:
            logger.error(f"Error initiating reprocessing of failed documents: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/cancel_pipeline",
        response_model=CancelPipelineResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def cancel_pipeline():
        """
        Request cancellation of the currently running pipeline.

        This endpoint sets a cancellation flag in the pipeline status. The pipeline will:
        1. Check this flag at key processing points
        2. Stop processing new documents
        3. Cancel all running document processing tasks
        4. Mark all PROCESSING documents as FAILED with reason "User cancelled"

        The cancellation is graceful and ensures data consistency. Documents that have
        completed processing will remain in PROCESSED status.

        Returns:
            CancelPipelineResponse: Response with status and message
                - status="cancellation_requested": Cancellation flag has been set
                - status="not_busy": Pipeline is not currently running

        Raises:
            HTTPException: If an error occurs while setting cancellation flag (500).
        """
        try:
            from lightrag.kg.shared_storage import (
                get_namespace_data,
                get_namespace_lock,
            )

            pipeline_status = await get_namespace_data(
                "pipeline_status", workspace=rag.workspace
            )
            pipeline_status_lock = get_namespace_lock(
                "pipeline_status", workspace=rag.workspace
            )

            async with pipeline_status_lock:
                if not pipeline_status.get("busy", False):
                    return CancelPipelineResponse(
                        status="not_busy",
                        message="Pipeline is not currently running. No cancellation needed.",
                    )

                # Set cancellation flag
                pipeline_status["cancellation_requested"] = True
                cancel_msg = "Pipeline cancellation requested by user"
                logger.info(cancel_msg)
                pipeline_status["latest_message"] = cancel_msg
                pipeline_status["history_messages"].append(cancel_msg)

            return CancelPipelineResponse(
                status="cancellation_requested",
                message="Pipeline cancellation has been requested. Documents will be marked as FAILED.",
            )

        except Exception as e:
            logger.error(f"Error requesting pipeline cancellation: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    return router
